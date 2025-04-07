import os
from pptx import Presentation
import re
import argparse
from collections import defaultdict

def sanitize_text(text):
    """VBA 문자열에 안전하게 사용할 수 있도록 텍스트 정리"""
    if text is None:
        return ""
    # 큰따옴표를 두 개로 이스케이프
    text = str(text).replace('"', '""')
    # 줄바꿈을 Chr(13)으로 변환
    text = text.replace('\n', '" & Chr(13) & "')
    return text

def get_shape_type(shape):
    """도형 유형 확인"""
    if shape.has_text_frame:
        return "TEXT"
    elif shape.has_table:
        return "TABLE"
    elif shape.has_chart:
        return "CHART"
    elif hasattr(shape, 'image'):
        return "PICTURE"
    elif shape.shape_type == 1:  # Auto-shape
        return "SHAPE"
    else:
        return "UNKNOWN"

def analyze_pptx(pptx_path):
    """PPTX 파일 분석"""
    prs = Presentation(pptx_path)
    presentation_data = {
        'slides': [],
        'slide_count': len(prs.slides),
        'slide_width': prs.slide_width,
        'slide_height': prs.slide_height
    }
    
    for i, slide in enumerate(prs.slides):
        slide_data = {
            'index': i + 1,
            'layout': slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else f"Layout {slide.slide_layout.slide_layout_id}",
            'shapes': []
        }
        
        # 슬라이드 내 도형 분석
        for shape in slide.shapes:
            shape_data = {
                'name': shape.name,
                'type': get_shape_type(shape),
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height
            }
            
            # 텍스트 내용 추출
            if shape.has_text_frame:
                text_content = ""
                for paragraph in shape.text_frame.paragraphs:
                    text_content += paragraph.text + "\n"
                shape_data['text'] = text_content.strip()
            
            # 테이블 내용 추출
            elif shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        cell_text = ""
                        for paragraph in cell.text_frame.paragraphs:
                            cell_text += paragraph.text + "\n"
                        row_data.append(cell_text.strip())
                    table_data.append(row_data)
                shape_data['table_data'] = table_data
                shape_data['rows'] = len(shape.table.rows)
                shape_data['columns'] = len(shape.table.columns)
            
            # 차트 데이터 기본 정보 추출 (상세 데이터 추출은 복잡함)
            elif shape.has_chart:
                shape_data['chart_type'] = shape.chart.chart_type
                
            slide_data['shapes'].append(shape_data)
        
        presentation_data['slides'].append(slide_data)
    
    return presentation_data

def generate_vba_code(presentation_data):
    """VBA 코드 생성"""
    vba_code = """Sub CreatePresentationFromAnalysis()
    ' PowerPoint 애플리케이션 객체 생성
    Dim ppApp As Object
    Dim ppPres As Object
    Dim ppSlide As Object
    Dim shp As Object
    
    ' PowerPoint 시작
    Set ppApp = CreateObject("PowerPoint.Application")
    ppApp.Visible = True ' PowerPoint 창을 보이게 설정
    
    ' 새 프레젠테이션 생성
    Set ppPres = ppApp.Presentations.Add
    
"""
    
    # 슬라이드 크기 설정
    vba_code += f"    ' 슬라이드 크기 설정\n"
    vba_code += f"    ppPres.PageSetup.SlideWidth = {presentation_data['slide_width']}\n"
    vba_code += f"    ppPres.PageSetup.SlideHeight = {presentation_data['slide_height']}\n\n"
    
    # 각 슬라이드에 대한 코드 생성
    for slide in presentation_data['slides']:
        slide_idx = slide['index']
        layout_map = {
            'Title Slide': 1,
            'Title and Content': 2,
            'Section Header': 3,
            'Two Content': 4,
            'Comparison': 5,
            'Title Only': 6,
            'Blank': 7,
            'Content with Caption': 8,
            'Picture with Caption': 9
        }
        
        # 레이아웃 번호 결정 (기본값: 빈 슬라이드)
        layout_num = 7  # 기본값: 빈 슬라이드
        for layout_name, layout_id in layout_map.items():
            if layout_name in slide['layout']:
                layout_num = layout_id
                break
        
        vba_code += f"    ' ----- 슬라이드 {slide_idx} 추가 ({slide['layout']}) -----\n"
        vba_code += f"    Set ppSlide = ppPres.Slides.Add({slide_idx}, {layout_num})\n\n"
        
        # 슬라이드 내 각 도형에 대한 코드 생성
        for i, shape in enumerate(slide['shapes']):
            shape_type = shape['type']
            
            if shape_type == "TEXT":
                text_content = sanitize_text(shape.get('text', ''))
                vba_code += f"    ' 텍스트 도형 추가: {shape['name']}\n"
                
                # 제목 도형인지 확인 (이름에 "Title"이 포함되어 있으면)
                if "title" in shape['name'].lower():
                    vba_code += f"    ppSlide.Shapes.Title.TextFrame.TextRange.Text = \"{text_content}\"\n\n"
                else:
                    vba_code += f"    Set shp = ppSlide.Shapes.AddTextbox(Orientation:=msoTextOrientationHorizontal, _\n"
                    vba_code += f"        Left:={shape['left']}, Top:={shape['top']}, _\n"
                    vba_code += f"        Width:={shape['width']}, Height:={shape['height']})\n"
                    vba_code += f"    shp.TextFrame.TextRange.Text = \"{text_content}\"\n\n"
            
            elif shape_type == "TABLE" and 'table_data' in shape:
                rows = shape.get('rows', 0)
                columns = shape.get('columns', 0)
                vba_code += f"    ' 표 추가 ({rows}행 {columns}열)\n"
                vba_code += f"    Set shp = ppSlide.Shapes.AddTable(NumRows:={rows}, NumColumns:={columns}, _\n"
                vba_code += f"        Left:={shape['left']}, Top:={shape['top']}, _\n"
                vba_code += f"        Width:={shape['width']}, Height:={shape['height']})\n"
                
                # 표 데이터 채우기
                for r, row in enumerate(shape['table_data']):
                    for c, cell_text in enumerate(row):
                        cell_text = sanitize_text(cell_text)
                        vba_code += f"    shp.Table.Cell({r+1}, {c+1}).Shape.TextFrame.TextRange.Text = \"{cell_text}\"\n"
                vba_code += "\n"
                
            elif shape_type == "CHART":
                vba_code += f"    ' 차트 추가 (차트 타입: {shape.get('chart_type', 'Unknown')})\n"
                vba_code += f"    Set shp = ppSlide.Shapes.AddChart(Type:=xlColumnClustered, _\n"
                vba_code += f"        Left:={shape['left']}, Top:={shape['top']}, _\n"
                vba_code += f"        Width:={shape['width']}, Height:={shape['height']})\n\n"
                
            elif shape_type == "PICTURE":
                vba_code += f"    ' 이미지 플레이스홀더 추가\n"
                vba_code += f"    ' 실제 이미지를 추가하려면 이미지 파일 경로 지정 필요\n"
                vba_code += f"    Set shp = ppSlide.Shapes.AddPicture( _\n"
                vba_code += f"        FileName:=\"C:\\Path\\To\\Image.jpg\", _\n"
                vba_code += f"        LinkToFile:=False, SaveWithDocument:=True, _\n"
                vba_code += f"        Left:={shape['left']}, Top:={shape['top']}, _\n"
                vba_code += f"        Width:={shape['width']}, Height:={shape['height']})\n\n"
                
            elif shape_type == "SHAPE":
                vba_code += f"    ' 기본 도형 추가\n"
                vba_code += f"    Set shp = ppSlide.Shapes.AddShape(Type:=msoShapeRectangle, _\n"
                vba_code += f"        Left:={shape['left']}, Top:={shape['top']}, _\n"
                vba_code += f"        Width:={shape['width']}, Height:={shape['height']})\n\n"
    
    # 프레젠테이션 저장 및 마무리
    vba_code += """    ' 프레젠테이션 저장
    ppPres.SaveAs "C:\\Path\\To\\ReconstructedPresentation.pptx"
    
    ' 객체 해제
    Set ppSlide = Nothing
    Set ppPres = Nothing
    Set ppApp = Nothing
    
    MsgBox "프레젠테이션이 성공적으로 생성되었습니다."
End Sub"""
    
    return vba_code

def main():
    parser = argparse.ArgumentParser(description='PPTX 파일을 분석하여 VBA 코드를 생성합니다.')
    parser.add_argument('pptx_file', help='분석할 PPTX 파일 경로')
    parser.add_argument('--output', '-o', help='VBA 코드를 저장할 파일 경로 (기본값: pptx 파일명에 _vba.bas 추가)')
    
    args = parser.parse_args()
    
    if not os.path.exists(args.pptx_file):
        print(f"오류: '{args.pptx_file}' 파일을 찾을 수 없습니다.")
        return
    
    # 출력 파일 경로 결정
    if args.output:
        output_file = args.output
    else:
        basename = os.path.splitext(args.pptx_file)[0]
        output_file = f"{basename}_vba.bas"
    
    try:
        print(f"'{args.pptx_file}' 파일 분석 중...")
        presentation_data = analyze_pptx(args.pptx_file)
        
        print(f"VBA 코드 생성 중...")
        vba_code = generate_vba_code(presentation_data)
        
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(vba_code)
        
        print(f"VBA 코드가 '{output_file}'에 저장되었습니다.")
    
    except Exception as e:
        print(f"오류 발생: {str(e)}")

if __name__ == "__main__":
    main()
    
# 실행: 
 # python .\pptx_to_vba.py  .\TeXBLEU_ppt.pptx