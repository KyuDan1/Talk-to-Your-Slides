
# 수정이 필요함.

from pptx import Presentation
from pptx.util import Emu
import os

def create_presentation():
    # 새 파워포인트 프레젠테이션 생성 및 슬라이드 크기 설정
    prs = Presentation()
    # VBA 코드에서 설정한 슬라이드 크기: 12192000 x 6858000 (단위: EMU)
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    # 기본적으로 사용 가능한 blank 슬라이드 레이아웃 선택
    blank_slide_layout = prs.slide_layouts[6]

    # ----- 슬라이드 1 (제목 슬라이드) 추가 -----
    slide1 = prs.slides.add_slide(blank_slide_layout)

    # 제목 텍스트 상자 추가
    title_box = slide1.shapes.add_textbox(Emu(1160584), Emu(884855), Emu(9870831), Emu(2387600))
    title_frame = title_box.text_frame
    title_frame.text = "TeXBLEU: Automatic Metric for Evaluate LaTeX Format"

    # 부제목 텍스트 상자 추가
    subtitle_box = slide1.shapes.add_textbox(Emu(1956079), Emu(3364568), Emu(8279842), Emu(1655762))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = (
        "Kyudan Jung1 , Nam-Joon Kim2∗ , Hyun Gon Ryu3 , Sieun Hyeon2 , Seung-jun Lee1 , Hyuk-jae Lee2\n"
        "1Chung-Ang University, Seoul, Korea\n"
        "2Seoul National Univeristy, Seoul, Korea\n"
        "3NVIDIA"
    )

    # 이미지 플레이스홀더 추가 (이미지 파일 경로를 실제 파일 경로로 변경)
    image_path = r"C:\Path\To\Image.jpg"  # 실제 이미지 파일 경로로 수정 필요
    slide1.shapes.add_picture(image_path, Emu(0), Emu(5270359), Emu(1969477), Emu(1587641))
    slide1.shapes.add_picture(image_path, Emu(203977), Emu(260510), Emu(1561521), Emu(890067))
    slide1.shapes.add_picture(image_path, Emu(10578549), Emu(170309), Emu(1409474), Emu(1042255))
    slide1.shapes.add_picture(image_path, Emu(9334500), Emu(5649971), Emu(2857500), Emu(1208029))

    # 이메일 텍스트 상자 추가
    email_box = slide1.shapes.add_textbox(Emu(4655828), Emu(5842391), Emu(2893741), Emu(369332))
    email_frame = email_box.text_frame
    email_frame.text = "wjdrbeks1021@gmail.com"

    # ----- 슬라이드 2 (Overview 슬라이드) 추가 -----
    slide2 = prs.slides.add_slide(blank_slide_layout)

    # 제목 텍스트 상자 추가
    title2_box = slide2.shapes.add_textbox(Emu(838200), Emu(365125), Emu(10515600), Emu(1325563))
    title2_frame = title2_box.text_frame
    title2_frame.text = "Overview"

    # 내용 텍스트 상자 추가
    content2_box = slide2.shapes.add_textbox(Emu(838200), Emu(1527349), Emu(10515600), Emu(4965526))
    content2_frame = content2_box.text_frame
    content2_frame.text = (
        "Problem Define\n"
        "1.1 Metric for LaTeX format\n"
        "1.2 Limitations of existing metrics\n"
        "Method\n"
        "2.1 Tokenizer trained by arXiv papers\n"
        "2.2 Fine-tuned embedding model\n"
        "2.3 Computing TeXBLEU algorithm\n"
        "Experiments\n"
        "3.1 Dataset and setup\n"
        "3.2 Human evaluation\n"
        "3.3 Result\n"
        "Conclusion"
    )

    # 슬라이드 번호 텍스트 상자 추가
    slide2_num = slide2.shapes.add_textbox(Emu(8610600), Emu(6356350), Emu(2743200), Emu(365125))
    slide2_num_frame = slide2_num.text_frame
    slide2_num_frame.text = "2/30"

    # 추가 슬라이드는 위와 같은 패턴으로 여러 개의 슬라이드를 생성할 수 있습니다.
    # 필요한 경우 슬라이드마다 텍스트 상자, 이미지, 도형 등을 추가하면 됩니다.

    # 프레젠테이션 저장 (파일 경로 업데이트 필요)
    save_path = r"C:\Path\To\ReconstructedPresentation.pptx"
    prs.save(save_path)
    print(f"Presentation saved as: {save_path}")

if __name__ == "__main__":
    create_presentation()
